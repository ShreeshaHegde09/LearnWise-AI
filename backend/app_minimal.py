from flask import Flask, request, jsonify
from flask_cors import CORS
from flask_sqlalchemy import SQLAlchemy
import os
from datetime import datetime
from werkzeug.utils import secure_filename
from gemini_service import GeminiService
# Note: attention_service removed - attention tracking now handled by frontend AttentionTracker component
from content_extractor import content_extractor
import json
import PyPDF2
import docx
from pptx import Presentation

app = Flask(__name__)
app.config['SECRET_KEY'] = 'your-secret-key'
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///learning_system.db'
app.config['UPLOAD_FOLDER'] = 'uploads'
CORS(app)

db = SQLAlchemy(app)

# Create uploads directory if it doesn't exist
os.makedirs(app.config['UPLOAD_FOLDER'], exist_ok=True)

# Initialize Gemini service
try:
    gemini_service = GeminiService()
    print("✓ Gemini API initialized successfully")
except Exception as e:
    print(f"Warning: Gemini API not available: {e}")
    gemini_service = None

# Emotion detection is handled client-side in the browser using TensorFlow.js
# Backend emotion service provides fallback support
try:
    from emotion_ensemble_service import emotion_service
    print("✓ Emotion ensemble service loaded (fallback mode)")
except Exception as e:
    print(f"Warning: Emotion service not available: {e}")
    emotion_service = None
print("ℹ️  Primary emotion detection runs client-side in browser (TensorFlow.js)")
print("   Backend provides fallback support if needed")

# Database Models
class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    created_at = db.Column(db.DateTime, default=datetime.utcnow)

class LearningMaterial(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    title = db.Column(db.String(200), nullable=False)  # Book/Document title
    topic = db.Column(db.String(200), nullable=True)  # FIXED: Made nullable
    filename = db.Column(db.String(200), nullable=False)
    content = db.Column(db.Text)  # Full content stored
    uploaded_at = db.Column(db.DateTime, default=datetime.utcnow)
    last_accessed = db.Column(db.DateTime, default=datetime.utcnow)

class LearningSession(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    material_id = db.Column(db.Integer, db.ForeignKey('learning_material.id'), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)
    topic = db.Column(db.String(200), nullable=False)  # Specific topic/chapter
    chunks = db.Column(db.Text)  # JSON string of generated chunks
    progress = db.Column(db.Integer, default=0)  # Percentage completed
    created_at = db.Column(db.DateTime, default=datetime.utcnow)
    updated_at = db.Column(db.DateTime, default=datetime.utcnow, onupdate=datetime.utcnow)

# Helper function to extract content from files
def extract_file_content(filepath, filename):
    """Extract text content from uploaded files"""
    content = ""
    file_ext = filename.lower().split('.')[-1]
    
    try:
        if file_ext == 'pdf':
            with open(filepath, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    content += page.extract_text() + "\n"
        
        elif file_ext in ['doc', 'docx']:
            doc = docx.Document(filepath)
            for paragraph in doc.paragraphs:
                content += paragraph.text + "\n"
        
        elif file_ext in ['ppt', 'pptx']:
            prs = Presentation(filepath)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        content += shape.text + "\n"
        
        elif file_ext == 'txt':
            with open(filepath, 'r', encoding='utf-8') as file:
                content = file.read()
    
    except Exception as e:
        print(f"Error extracting content from {filename}: {e}")
        content = ""
    
    return content.strip()

# Routes
@app.route('/api/health', methods=['GET'])
def health_check():
    return jsonify({"status": "healthy", "message": "Backend is running!"})

@app.route('/api/upload-material', methods=['POST'])
def upload_material():
    try:
        if 'file' not in request.files:
            return jsonify({"error": "No file provided"}), 400
        
        file = request.files['file']
        topic = request.form.get('topic')
        title = request.form.get('title', file.filename)  # Book/document title
        user_id = request.form.get('user_id', 1)
        
        if file.filename == '':
            return jsonify({"error": "No file selected"}), 400
        
        if file and topic:
            filename = secure_filename(file.filename)
            filepath = os.path.join(app.config['UPLOAD_FOLDER'], filename)
            file.save(filepath)
            
            # Extract FULL content from file
            print(f"📄 Extracting content from {filename}...")
            full_content = extract_file_content(filepath, filename)
            print(f"📝 Extracted {len(full_content)} characters")
            
            if not full_content or len(full_content.strip()) < 50:
                print(f"⚠️  Content too short or empty, using fallback")
                full_content = f"Learning material about {topic}. This file contains important information about {topic} that we'll explore together."
            else:
                print(f"✅ Content extracted successfully")
            
            # Check if this material already exists
            existing_material = LearningMaterial.query.filter_by(
                user_id=user_id,
                filename=filename
            ).first()
            
            if existing_material:
                # Update existing material
                print(f"📚 Material already exists, updating...")
                existing_material.content = full_content
                existing_material.topic = topic  # Update topic as well
                existing_material.last_accessed = datetime.utcnow()
                material = existing_material
            else:
                # Save NEW material to database with FULL content
                material = LearningMaterial(
                    user_id=user_id,
                    title=title,
                    topic=topic,  # FIXED: Add topic field
                    filename=filename,
                    content=full_content  # Store full content
                )
                db.session.add(material)
            
            db.session.commit()
            print(f"💾 Material saved with ID: {material.id}")
            
            # SMART EXTRACTION: Extract only topic-relevant content
            print(f"🔍 Extracting topic-relevant content for: {topic}")
            relevant_content = content_extractor.extract_topic_content(full_content, topic)
            print(f"✂️  Extracted {len(relevant_content)} characters (from {len(full_content)})")
            print(f"📊 Reduction: {100 - int(len(relevant_content)/len(full_content)*100)}%")
            
            # Generate learning chunks using ONLY relevant content
            chunks = []
            if gemini_service:
                try:
                    print(f"🤖 Generating chunks with Gemini for topic: {topic}")
                    chunks = gemini_service.generate_learning_chunks(relevant_content, topic)
                    print(f"✅ Generated {len(chunks)} chunks for topic: {topic}")
                    if chunks:
                        print(f"📌 First chunk title: {chunks[0].get('title', 'N/A')}")
                except Exception as e:
                    print(f"❌ Error generating chunks: {e}")
                    import traceback
                    traceback.print_exc()
                    chunks = []
            else:
                print(f"⚠️  Gemini service not available")
            
            # Fallback chunks if Gemini fails or no chunks generated
            if not chunks or len(chunks) == 0:
                print(f"⚠️  Using fallback chunks (Gemini returned {len(chunks) if chunks else 0} chunks)")
                
                # Extract more content for fallback
                content_parts = content.split('\n\n') if content else []
                content_preview = content[:1000] if content else f"Learning material about {topic}"
                
                # Try to create more meaningful fallback based on actual content
                chunks = [
                    {
                        "id": 1,
                        "title": f"🎯 Introduction to {topic}",
                        "content": f"""
                        <h3>🎯 Understanding {topic}</h3>
                        <p>Welcome to your comprehensive guide on <strong>{topic}</strong>.</p>
                        <p><em>This material covers essential concepts and practical applications.</em></p>
                        <h3>📚 Content Overview</h3>
                        <p>{content_preview[:400] if len(content_preview) > 400 else content_preview}</p>
                        <h3>🔑 What You'll Master</h3>
                        <ul>
                        <li>Fundamental concepts and definitions</li>
                        <li>Core principles and methodologies</li>
                        <li>Practical examples and use cases</li>
                        <li>Advanced techniques and best practices</li>
                        </ul>
                        """,
                        "estimated_time": "8 min",
                        "objectives": ["Understand core concepts", "Learn fundamental principles", "Grasp key terminology"]
                    },
                    {
                        "id": 2,
                        "title": f"📖 Deep Dive into {topic}",
                        "content": f"""
                        <h3>📖 Core Concepts and Details</h3>
                        <p>Let's explore the <strong>detailed aspects</strong> of {topic}.</p>
                        <p>{content_preview[400:800] if len(content_preview) > 800 else content_preview[400:] if len(content_preview) > 400 else "This section covers the main concepts in detail."}</p>
                        <h3>💡 Key Principles</h3>
                        <ul>
                        <li><strong>Foundation:</strong> Building blocks and basic structure</li>
                        <li><strong>Methodology:</strong> How to approach and apply concepts</li>
                        <li><strong>Techniques:</strong> Specific methods and approaches</li>
                        <li><strong>Standards:</strong> Best practices and conventions</li>
                        </ul>
                        <p><em>Understanding these principles is crucial for mastery.</em></p>
                        """,
                        "estimated_time": "12 min",
                        "objectives": ["Master detailed concepts", "Understand methodologies", "Learn techniques"]
                    },
                    {
                        "id": 3,
                        "title": f"⚡ Advanced Topics in {topic}",
                        "content": f"""
                        <h3>⚡ Advanced Concepts</h3>
                        <p>Now we'll cover <strong>advanced aspects</strong> and deeper understanding of {topic}.</p>
                        <p>{content_preview[800:1200] if len(content_preview) > 1200 else content_preview[800:] if len(content_preview) > 800 else "Advanced topics build on the fundamentals."}</p>
                        <h3>🚀 Advanced Techniques</h3>
                        <ul>
                        <li>Complex scenarios and solutions</li>
                        <li>Optimization and efficiency</li>
                        <li>Edge cases and special considerations</li>
                        <li>Integration with other concepts</li>
                        </ul>
                        """,
                        "estimated_time": "10 min",
                        "objectives": ["Master advanced concepts", "Handle complex scenarios", "Optimize solutions"]
                    },
                    {
                        "id": 4,
                        "title": f"💼 Practical Applications",
                        "content": f"""
                        <h3>💼 Real-World Applications</h3>
                        <p>Let's see how <strong>{topic}</strong> is applied in real-world scenarios.</p>
                        <h3>✨ Use Cases</h3>
                        <ul>
                        <li><strong>Industry Applications:</strong> How professionals use this daily</li>
                        <li><strong>Common Patterns:</strong> Frequently used approaches</li>
                        <li><strong>Problem Solving:</strong> Applying concepts to solve real problems</li>
                        <li><strong>Case Studies:</strong> Examples from actual projects</li>
                        </ul>
                        <p><em>These practical examples help solidify your understanding.</em></p>
                        <h3>🎯 Implementation Tips</h3>
                        <ul>
                        <li>Start with simple examples</li>
                        <li>Build complexity gradually</li>
                        <li>Test your understanding with practice</li>
                        <li>Apply to your own projects</li>
                        </ul>
                        """,
                        "estimated_time": "10 min",
                        "objectives": ["Apply knowledge practically", "Understand real use cases", "Solve real problems"]
                    },
                    {
                        "id": 5,
                        "title": f"🎓 Mastery and Next Steps",
                        "content": f"""
                        <h3>🎓 Achieving Mastery</h3>
                        <p>You've covered the essentials of <strong>{topic}</strong>. Now let's consolidate your knowledge.</p>
                        <h3>📝 Key Takeaways</h3>
                        <ul>
                        <li>Core concepts and their relationships</li>
                        <li>Practical applications and use cases</li>
                        <li>Best practices and common patterns</li>
                        <li>Advanced techniques for complex scenarios</li>
                        </ul>
                        <h3>🚀 Continue Your Learning</h3>
                        <p><em>To truly master {topic}, consider:</em></p>
                        <ul>
                        <li><strong>Practice:</strong> Apply concepts to real projects</li>
                        <li><strong>Explore:</strong> Dive deeper into advanced topics</li>
                        <li><strong>Build:</strong> Create your own examples and solutions</li>
                        <li><strong>Share:</strong> Teach others to reinforce your knowledge</li>
                        </ul>
                        <p>Remember: Mastery comes from consistent practice and application!</p>
                        """,
                        "estimated_time": "8 min",
                        "objectives": ["Consolidate knowledge", "Plan next steps", "Build confidence"]
                    }
                ]
            
            # Save learning session
            session = LearningSession(
                material_id=material.id,
                user_id=user_id,
                topic=topic,
                chunks=json.dumps(chunks),  # Store chunks as JSON
                progress=0
            )
            db.session.add(session)
            db.session.commit()
            
            response_data = {
                "message": "File uploaded successfully",
                "material_id": material.id,
                "session_id": session.id,
                "topic": topic,
                "title": title,
                "filename": filename,
                "content_preview": relevant_content[:200] if relevant_content else "",
                "chunks": chunks,
                "is_new_upload": existing_material is None
            }
            
            print(f"📤 Sending response with {len(chunks)} chunks")
            print(f"📋 Session ID: {session.id}, Material ID: {material.id}")
            
            return jsonify(response_data)
        
        return jsonify({"error": "Missing file or topic"}), 400
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/materials/<int:user_id>', methods=['GET'])
def get_user_materials(user_id):
    """Get all uploaded materials for a user"""
    materials = LearningMaterial.query.filter_by(user_id=user_id).order_by(LearningMaterial.last_accessed.desc()).all()
    return jsonify([{
        "id": m.id,
        "title": m.title,
        "filename": m.filename,
        "uploaded_at": m.uploaded_at.isoformat(),
        "last_accessed": m.last_accessed.isoformat(),
        "content_length": len(m.content),
        "content_preview": m.content[:200] + "..." if len(m.content) > 200 else m.content
    } for m in materials])

@app.route('/api/materials/<int:material_id>/sessions', methods=['GET'])
def get_material_sessions(material_id):
    """Get all learning sessions for a material"""
    sessions = LearningSession.query.filter_by(material_id=material_id).order_by(LearningSession.updated_at.desc()).all()
    return jsonify([{
        "id": s.id,
        "topic": s.topic,
        "progress": s.progress,
        "created_at": s.created_at.isoformat(),
        "updated_at": s.updated_at.isoformat()
    } for s in sessions])

@app.route('/api/continue-learning', methods=['POST'])
def continue_learning():
    """Continue learning from an existing material with a new topic"""
    try:
        data = request.json
        material_id = data.get('material_id')  # This is Firestore ID (string)
        topic = data.get('topic')
        user_id = data.get('user_id', 1)
        material_content = data.get('content')  # Get content from frontend
        material_title = data.get('title', 'Study Material')
        
        print(f"📥 Continue learning request: material_id={material_id}, topic={topic}")
        
        if not topic:
            return jsonify({"error": "topic required"}), 400
        
        if not material_content:
            return jsonify({"error": "material content required"}), 400
        
        print(f"📚 Continuing learning from material: {material_title}")
        print(f"🎯 New topic: {topic}")
        print(f"📄 Content length: {len(material_content)} characters")
        
        # SMART EXTRACTION: Extract only topic-relevant content
        print(f"🔍 Extracting topic-relevant content for: {topic}")
        relevant_content = content_extractor.extract_topic_content(material_content, topic)
        print(f"✂️  Extracted {len(relevant_content)} characters (from {len(material_content)})")
        print(f"📊 Reduction: {100 - int(len(relevant_content)/len(material_content)*100)}%")
        
        # Generate learning chunks using ONLY relevant content
        chunks = []
        if gemini_service:
            try:
                print(f"🤖 Generating chunks with Gemini for topic: {topic}")
                chunks = gemini_service.generate_learning_chunks(relevant_content, topic)
                print(f"✅ Generated {len(chunks)} chunks")
            except Exception as e:
                print(f"❌ Error generating chunks: {e}")
                chunks = []
        
        # Fallback chunks if needed
        if not chunks or len(chunks) == 0:
            print(f"⚠️  Using fallback chunks")
            chunks = generate_fallback_chunks(topic, material_title, relevant_content)
        
        # Return the generated chunks (no need to save to SQLite since using Firestore)
        return jsonify({
            "message": "New learning content generated successfully",
            "material_id": material_id,  # Firestore ID
            "session_id": f"session_{int(datetime.utcnow().timestamp())}",
            "topic": topic,
            "title": material_title,
            "chunks": chunks,
            "is_continuation": True
        })
        
    except Exception as e:
        print(f"Error in continue_learning: {e}")
        import traceback
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

def generate_fallback_chunks(topic, filename, content):
    """Generate fallback chunks when Gemini fails"""
    content_parts = content.split('\n\n') if content else []
    content_preview = content[:1000] if content else f"Learning material about {topic}"
    
    return [
        {
            "id": 1,
            "title": f"🎯 Introduction to {topic}",
            "content": f"""
            <h3>🎯 Understanding {topic}</h3>
            <p>Welcome to your comprehensive guide on <strong>{topic}</strong>.</p>
            <p><em>This material covers essential concepts and practical applications.</em></p>
            <h3>📚 Content Overview</h3>
            <p>{content_preview[:400] if len(content_preview) > 400 else content_preview}</p>
            <h3>🔑 What You'll Master</h3>
            <ul>
            <li>Fundamental concepts and definitions</li>
            <li>Core principles and methodologies</li>
            <li>Practical examples and use cases</li>
            <li>Advanced techniques and best practices</li>
            </ul>
            """,
            "estimated_time": "8 min"
        },
        {
            "id": 2,
            "title": f"📖 Deep Dive into {topic}",
            "content": f"""
            <h3>📖 Core Concepts and Details</h3>
            <p>Let's explore the <strong>detailed aspects</strong> of {topic}.</p>
            <p>{content_preview[400:800] if len(content_preview) > 800 else content_preview[400:] if len(content_preview) > 400 else "This section covers the main concepts in detail."}</p>
            <h3>💡 Key Principles</h3>
            <ul>
            <li><strong>Foundation:</strong> Building blocks and basic structure</li>
            <li><strong>Methodology:</strong> How to approach and apply concepts</li>
            <li><strong>Techniques:</strong> Specific methods and approaches</li>
            <li><strong>Standards:</strong> Best practices and conventions</li>
            </ul>
            """,
            "estimated_time": "12 min"
        }
    ]

@app.route('/api/generate-content', methods=['POST'])
def generate_content():
    # Mock response for testing
    data = request.json
    material_id = data.get('material_id')
    content_type = data.get('type', 'summary')
    
    material = LearningMaterial.query.get(material_id)
    if not material:
        return jsonify({"error": "Material not found"}), 404
    
    # Mock responses
    mock_responses = {
        'summary': f"Summary of {material.topic}: This is a mock summary of the uploaded material.",
        'quiz': {
            "questions": [
                {
                    "question": f"What is the main topic of this material about {material.topic}?",
                    "options": ["Option A", "Option B", "Option C", "Option D"],
                    "correct": 0
                }
            ]
        },
        'flashcards': [
            {
                "front": f"Key concept in {material.topic}",
                "back": "This is a mock flashcard content"
            }
        ]
    }
    
    return jsonify({
        "content": mock_responses.get(content_type, "Mock content generated"),
        "type": content_type
    })

@app.route('/api/simplify-content', methods=['POST'])
def simplify_content():
    try:
        data = request.json
        content = data.get('content')
        
        if not gemini_service:
            return jsonify({"error": "AI service not available"}), 503
        
        simplified = gemini_service.simplify_content(content)
        
        return jsonify({
            "simplified_content": simplified
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/generate-quiz', methods=['POST'])
def generate_quiz():
    try:
        data = request.json
        
        # Support both material_id and direct topic/content
        material_id = data.get('material_id')
        topic = data.get('topic')
        content = data.get('content')
        question_count = data.get('question_count', 5)
        
        if material_id:
            material = LearningMaterial.query.get(material_id)
            if not material:
                return jsonify({"error": "Material not found"}), 404
            topic = material.topic
            content = material.content
        elif not topic or not content:
            return jsonify({"error": "Topic and content required"}), 400
        
        if not gemini_service:
            return jsonify({"error": "AI service not available"}), 503
        
        quiz = gemini_service.generate_quiz(content, topic, question_count)
        
        return jsonify(quiz)
        
    except Exception as e:
        print(f"Error in generate_quiz: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/chat', methods=['POST'])
def chat():
    try:
        data = request.json
        message = data.get('message')
        current_chunk = data.get('current_chunk', {})
        context = data.get('context', '')
        
        if not gemini_service:
            return jsonify({"error": "AI service not available"}), 503
        
        # Build comprehensive context from current chunk
        full_context = f"""
        Current learning topic: {current_chunk.get('title', 'Unknown')}
        Content: {current_chunk.get('content', '')}
        Context: {context}
        """
        
        response = gemini_service.chat_response(message, full_context)
        
        return jsonify({
            "response": response
        })
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/screen-data', methods=['POST', 'OPTIONS'])
def handle_screen_data():
    if request.method == 'OPTIONS':
        return '', 200
    
    try:
        data = request.json
        session_id = data.get('session_id')
        screen_data = data.get('screen_data')
        
        # For now, just log the screen data
        print(f"Screen data received for session {session_id}: {len(screen_data) if screen_data else 0} bytes")
        
        return jsonify({"status": "success", "message": "Screen data received"})
        
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/api/generate-flashcards', methods=['POST'])
def generate_flashcards():
    try:
        data = request.json
        topic = data.get('topic')
        content = data.get('content')
        card_count = data.get('card_count', 10)
        
        if not topic or not content:
            return jsonify({"error": "Topic and content required"}), 400
        
        if not gemini_service:
            return jsonify({"error": "AI service not available"}), 503
        
        flashcards = gemini_service.generate_flashcards(content, topic, card_count)
        
        return jsonify(flashcards)
        
    except Exception as e:
        print(f"Error in generate_flashcards: {e}")
        return jsonify({"error": str(e)}), 500

@app.route('/api/emotion-data', methods=['POST'])
def handle_emotion_data():
    try:
        # Handle both JSON and form data
        if request.is_json:
            data = request.json or {}
        else:
            data = request.form.to_dict()
        
        session_id = data.get('session_id', 'unknown')
        emotion_data = data.get('emotion_data', 'no_data')
        
        # For now, just log the emotion data
        print(f"Emotion data received for session {session_id}: {emotion_data}")
        
        # You could store this in the database or trigger actions based on emotions
        # For example, if stress is detected, suggest a break or simplify content
        
        return jsonify({"status": "success", "message": "Emotion data received"})
        
    except Exception as e:
        print(f"Error handling emotion data: {e}")
        return jsonify({"status": "success", "message": "Emotion data received (with errors)"}), 200

# ============= ATTENTION TRACKING ENDPOINTS =============
# NOTE: These endpoints are commented out because attention tracking is now
# handled entirely by the frontend AttentionTracker React component.
# See: NovProject/frontend/src/components/AttentionTracker.tsx
# Documentation: .kiro/specs/attention-tracking-system/DEVELOPER_GUIDE.md

# @app.route('/api/attention/event', methods=['POST', 'OPTIONS'])
# def handle_attention_event():
#     """Receive attention tracking events from Chrome extension"""
#     if request.method == 'OPTIONS':
#         return '', 200
#     
#     try:
#         data = request.json
#         
#         # Add event to attention tracker
#         attention_tracker.add_event(data)
#         
#         # Calculate current attention
#         attention_status = attention_tracker.calculate_attention_score()
#         
#         # Check if we should show an alert
#         should_alert, alert_message, alert_type = attention_tracker.should_show_alert()
#         
#         response = {
#             "status": "success",
#             "attention": attention_status,
#             "show_alert": should_alert,
#             "alert_message": alert_message,
#             "alert_type": alert_type
#         }
#         
#         # Log for debugging
#         if should_alert:
#             print(f"⚠️  Low attention detected: {alert_message}")
#         
#         return jsonify(response)
#         
#     except Exception as e:
#         print(f"Error handling attention event: {e}")
#         return jsonify({"status": "error", "message": str(e)}), 500

# @app.route('/api/attention/status', methods=['GET'])
# def get_attention_status():
#     """Get current attention status"""
#     try:
#         attention_status = attention_tracker.calculate_attention_score()
#         summary = attention_tracker.get_attention_summary()
#         should_alert, alert_message, alert_type = attention_tracker.should_show_alert()
#         
#         return jsonify({
#             "status": "success",
#             "current_attention": attention_status,
#             "summary": summary,
#             "show_alert": should_alert,
#             "alert_message": alert_message,
#             "alert_type": alert_type
#         })
#         
#     except Exception as e:
#         return jsonify({"status": "error", "message": str(e)}), 500

# @app.route('/api/attention/reset', methods=['POST'])
# def reset_attention_tracking():
#     """Reset attention tracking data"""
#     try:
#         attention_tracker.reset()
#         return jsonify({"status": "success", "message": "Attention tracking reset"})
#     except Exception as e:
#         return jsonify({"status": "error", "message": str(e)}), 500

# @app.route('/api/attention/debug', methods=['GET'])
# def debug_attention():
#     """Debug endpoint to check attention tracking state"""
#     try:
#         return jsonify({
#             "status": "success",
#             "tab_left_time": attention_tracker.tab_left_time.isoformat() if attention_tracker.tab_left_time else None,
#             "away_alert_sent": attention_tracker.away_alert_sent,
#             "recent_events_count": len(attention_tracker.recent_events),
#             "attention_history_count": len(attention_tracker.attention_history),
#             "last_alert": attention_tracker.last_low_attention_alert.isoformat() if attention_tracker.last_low_attention_alert else None
#         })
#     except Exception as e:
#         return jsonify({"status": "error", "message": str(e)}), 500


# ============= EMOTION DETECTION ENDPOINTS =============
# Ensemble of 3 models: MobileNet, EfficientNet, LandmarkCNN

@app.route('/api/emotion/predict', methods=['POST'])
def predict_emotion():
    """
    Emotion detection endpoint - Returns info that detection happens client-side
    Note: Actual emotion detection runs in the browser using TensorFlow.js
    """
    # Emotion detection is handled client-side, not on backend
    return jsonify({
        'info': 'Emotion detection runs client-side in browser',
        'message': 'This endpoint is not used. Emotion detection happens in the browser using TensorFlow.js',
        'emotion': 'neutral',
        'confidence': 0.0,
        'client_side': True
    }), 200
    
    try:
        data = request.json
        
        # Validate request
        if not data or 'image' not in data:
            return jsonify({
                'error': 'Invalid request',
                'message': 'Missing required field: image'
            }), 400
        
        image_base64 = data.get('image')
        landmarks = data.get('landmarks')  # Optional (not used by current models)
        session_id = data.get('session_id')
        
        # Run ensemble prediction
        result = emotion_service.predict(image_base64, landmarks)
        
        # Add session_id to result if provided
        if session_id:
            result['session_id'] = session_id
        
        # Add timestamp
        result['timestamp'] = datetime.utcnow().isoformat()
        
        return jsonify(result), 200
        
    except ValueError as e:
        return jsonify({
            'error': 'Invalid input',
            'message': str(e)
        }), 400
    except RuntimeError as e:
        return jsonify({
            'error': 'Prediction failed',
            'message': str(e)
        }), 500
    except Exception as e:
        return jsonify({
            'error': 'Internal server error',
            'message': str(e)
        }), 500


@app.route('/api/emotion/predict/batch', methods=['POST'])
def predict_emotion_batch():
    """
    Predict emotions for multiple images in batch using ensemble
    """
    if not emotion_service or not emotion_service.models_loaded:
        return jsonify({
            'error': 'Emotion service not available',
            'message': 'Models not loaded on server'
        }), 503
    
    try:
        data = request.json
        
        # Validate request
        if not data or 'images' not in data:
            return jsonify({
                'error': 'Invalid request',
                'message': 'Missing required field: images'
            }), 400
        
        images = data.get('images', [])
        local_predictions = data.get('local_predictions')
        session_id = data.get('session_id')
        
        if not images:
            return jsonify({
                'error': 'Invalid request',
                'message': 'Images array is empty'
            }), 400
        
        # Run batch prediction
        results = emotion_service.predict_batch(images, local_predictions)
        
        return jsonify({
            'session_id': session_id,
            'predictions': results,
            'count': len(results),
            'timestamp': datetime.utcnow().isoformat()
        }), 200
        
    except Exception as e:
        return jsonify({
            'error': 'Internal server error',
            'message': str(e)
        }), 500


@app.route('/api/emotion/info', methods=['GET'])
def get_emotion_info():
    """
    Get information about the emotion detection service
    """
    if not emotion_service:
        return jsonify({
            'available': False,
            'message': 'Emotion service not initialized'
        }), 503
    
    info = emotion_service.get_model_info()
    info['available'] = True
    
    return jsonify(info), 200


@app.route('/api/generate-summary', methods=['POST'])
def generate_summary():
    """Generate a comprehensive summary for a specific topic"""
    try:
        data = request.json
        topic = data.get('topic', '')
        content = data.get('content', '')
        
        if not topic:
            return jsonify({'error': 'Topic is required'}), 400
        
        if not gemini_service:
            # Fallback summary if Gemini not available
            return jsonify({
                'summary': f"""Summary of {topic}:

This topic covers important concepts that you've been studying. Here are the key points:

• Core concepts and fundamental principles
• Practical applications and real-world examples
• Important relationships between different ideas
• Best practices and common approaches
• Areas for further exploration and mastery

Review your materials, flashcards, and quiz results to reinforce your understanding of these concepts.""",
                'topic': topic
            })
        
        # Generate comprehensive summary using Gemini
        if content:
            prompt = f"""Based on the following study materials about {topic}, generate a comprehensive summary:

{content[:3000]}

Please provide:
1. Key concepts and definitions
2. Main points to remember
3. Important relationships or connections
4. Practical applications
5. Common misconceptions to avoid

Keep it concise but informative, around 200-300 words."""
        else:
            prompt = f"""Generate a comprehensive summary for the topic: {topic}

Please provide:
1. Key concepts and definitions
2. Main points to remember
3. Important relationships or connections
4. Practical applications
5. Common misconceptions to avoid

Keep it concise but informative, around 200-300 words."""

        try:
            # Use the GeminiService instance that's already initialized
            summary = gemini_service.generate_summary(prompt)
            
            return jsonify({
                'summary': summary,
                'topic': topic
            })
        except Exception as e:
            print(f"Error generating summary with Gemini: {e}")
            # Fallback
            return jsonify({
                'summary': f"""Summary of {topic}:

This topic encompasses several important concepts that form the foundation of your learning. The key areas include fundamental principles, practical applications, and advanced techniques.

Understanding these concepts will help you:
• Master the core principles
• Apply knowledge in real-world scenarios
• Build on this foundation for advanced topics
• Avoid common pitfalls and misconceptions

Continue practicing with flashcards and quizzes to solidify your understanding.""",
                'topic': topic
            })
            
    except Exception as e:
        print(f"Error in generate_summary: {e}")
        return jsonify({'error': str(e)}), 500

if __name__ == '__main__':
    with app.app_context():
        # Try to migrate database schema
        try:
            from sqlalchemy import inspect, text
            inspector = inspect(db.engine)
            
            # Check if learning_material has the new columns
            columns = [col['name'] for col in inspector.get_columns('learning_material')]
            
            if 'title' not in columns:
                print("🔄 Migrating database schema...")
                # Add new columns to learning_material
                with db.engine.connect() as conn:
                    conn.execute(text('ALTER TABLE learning_material ADD COLUMN title VARCHAR(200)'))
                    conn.execute(text('ALTER TABLE learning_material ADD COLUMN last_accessed DATETIME'))
                    conn.commit()
                print("✅ Added title and last_accessed columns")
            
            # Check if learning_session table exists
            if 'learning_session' not in inspector.get_table_names():
                print("🔄 Creating learning_session table...")
                db.create_all()
                print("✅ Created learning_session table")
            
            # Update existing records
            with db.engine.connect() as conn:
                # Set title = filename for existing records
                conn.execute(text('UPDATE learning_material SET title = filename WHERE title IS NULL'))
                # Set last_accessed = uploaded_at for existing records
                conn.execute(text('UPDATE learning_material SET last_accessed = uploaded_at WHERE last_accessed IS NULL'))
                conn.commit()
            
            print("✅ Database schema is up to date!")
            
        except Exception as e:
            print(f"⚠️  Migration note: {e}")
            # If migration fails, just create all tables (for new databases)
            db.create_all()
    
    app.run(debug=True, port=5000)
    
    